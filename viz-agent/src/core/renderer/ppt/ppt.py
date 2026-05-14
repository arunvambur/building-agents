import io
import logging
import os
import tempfile
import uuid
from datetime import date
from typing import Any

import matplotlib
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from core.dsl.schema import Chart, VisualizationSpec
from core.renderer.base import Renderer
from core.renderer.image.image import ImageRenderer

matplotlib.use("Agg")

logger = logging.getLogger(__name__)

# Brand colours — kept in sync with image.py subtle palette
_BRAND_BLUE = RGBColor(0x7B, 0x9C, 0xCC)
_DARK       = RGBColor(0x1A, 0x1F, 0x2E)
_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
_MID_GRAY   = RGBColor(0x7A, 0x84, 0x99)

# Slide dimensions (widescreen 16:9)
_SLIDE_W = Cm(33.87)
_SLIDE_H = Cm(19.05)


class PPTRenderer(Renderer):

    name = "ppt"

    def __init__(self):
        self._image_renderer = ImageRenderer()

    def supports(self, format: str) -> bool:
        return format.lower() == "ppt"

    def render(self, spec: VisualizationSpec, data: Any) -> str:
        """
        Renders a VisualizationSpec to a PowerPoint (.pptx) file.
        Slide 1: title slide with report summary.
        Slides 2+: one slide per chart with embedded PNG and data table.
        Returns a 'file://<path>' string.
        """
        rows: list[dict] = data if isinstance(data, list) else []
        output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pptx")

        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H

        # Slide 1 — title
        self._add_title_slide(prs, spec, rows)

        # One slide per chart
        for i, chart_spec in enumerate(spec.charts):
            self._add_chart_slide(prs, chart_spec, rows, i + 1)

        # Filters slide (if any)
        if spec.filters:
            self._add_filters_slide(prs, spec)

        prs.save(output_path)
        logger.info("[ppt_renderer] written to %s", output_path)
        return f"file://{output_path}"

    # ---- slide builders ----

    def _add_title_slide(self, prs: Presentation, spec: VisualizationSpec, rows: list[dict]) -> None:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        self._fill_background(slide, _DARK)

        # Title
        self._add_text_box(
            slide,
            text="Viz Agent — Visualization Report",
            left=Cm(2), top=Cm(6), width=Cm(30), height=Cm(2.5),
            font_size=Pt(32), bold=True, color=_WHITE, align=PP_ALIGN.CENTER,
        )
        # Subtitle
        subtitle = (
            f"{len(spec.charts)} chart(s)  |  {len(rows)} records  |  "
            f"Generated {date.today().strftime('%d %b %Y')}"
        )
        self._add_text_box(
            slide,
            text=subtitle,
            left=Cm(2), top=Cm(9), width=Cm(30), height=Cm(1.2),
            font_size=Pt(14), bold=False, color=_MID_GRAY, align=PP_ALIGN.CENTER,
        )
        # Brand accent bar
        self._add_rect(slide, left=Cm(12), top=Cm(11.5), width=Cm(10), height=Cm(0.15), color=_BRAND_BLUE)

    def _add_chart_slide(
        self, prs: Presentation, chart_spec: Chart, rows: list[dict], index: int
    ) -> None:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        self._fill_background(slide, _LIGHT_GRAY)

        # Header bar
        self._add_rect(slide, left=Cm(0), top=Cm(0), width=_SLIDE_W, height=Cm(1.8), color=_BRAND_BLUE)

        # Slide title in header
        title = chart_spec.title or f"Chart {index}"
        self._add_text_box(
            slide,
            text=title,
            left=Cm(0.8), top=Cm(0.2), width=Cm(28), height=Cm(1.4),
            font_size=Pt(18), bold=True, color=_WHITE,
        )

        # Chart image (left side)
        chart_bytes = self._render_chart_to_bytes(chart_spec, rows)
        chart_stream = io.BytesIO(chart_bytes)
        slide.shapes.add_picture(chart_stream, left=Cm(0.5), top=Cm(2.2), width=Cm(20), height=Cm(12))

        # Data table (right side) — show top 10 rows to fit slide
        if rows:
            self._add_data_table(slide, rows[:10], left=Cm(21.5), top=Cm(2.2), width=Cm(11.5))

        # Footer
        self._add_text_box(
            slide,
            text=f"Slide {index + 1}  |  Viz Agent",
            left=Cm(0), top=Cm(18.2), width=_SLIDE_W, height=Cm(0.7),
            font_size=Pt(8), bold=False, color=_MID_GRAY, align=PP_ALIGN.CENTER,
        )

    def _add_filters_slide(self, prs: Presentation, spec: VisualizationSpec) -> None:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._fill_background(slide, _LIGHT_GRAY)

        self._add_rect(slide, left=Cm(0), top=Cm(0), width=_SLIDE_W, height=Cm(1.8), color=_BRAND_BLUE)
        self._add_text_box(
            slide,
            text="Filters Applied",
            left=Cm(0.8), top=Cm(0.2), width=Cm(28), height=Cm(1.4),
            font_size=Pt(18), bold=True, color=_WHITE,
        )

        y = Cm(2.5)
        for f in spec.filters:
            self._add_text_box(
                slide,
                text=f"{f.field}  {f.op}  {f.value}",
                left=Cm(2), top=y, width=Cm(28), height=Cm(0.9),
                font_size=Pt(14), bold=False, color=_DARK,
            )
            y += Cm(1.1)

    # ---- helpers ----

    def _render_chart_to_bytes(self, chart_spec: Chart, rows: list[dict]) -> bytes:
        fig, ax = plt.subplots(1, 1, figsize=(10, 6))
        self._image_renderer._apply_theme(fig)
        self._image_renderer._draw_chart(ax, chart_spec, rows)
        fig.tight_layout(pad=2.0)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    def _add_data_table(
        self, slide, rows: list[dict], left, top, width
    ) -> None:
        if not rows:
            return

        headers = list(rows[0].keys())
        n_cols = len(headers)
        n_rows = len(rows) + 1  # +1 for header

        row_h = Cm(0.7)
        col_w = width // n_cols
        table_h = row_h * n_rows

        table = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h).table

        # Style header row
        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = _BRAND_BLUE
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.bold = True
            run.font.size = Pt(8)
            run.font.color.rgb = _WHITE
            para.alignment = PP_ALIGN.CENTER

        # Data rows
        for ri, row in enumerate(rows):
            bg = _WHITE if ri % 2 == 0 else _LIGHT_GRAY
            for ci, header in enumerate(headers):
                cell = table.cell(ri + 1, ci)
                cell.text = str(row.get(header, ""))
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                para = cell.text_frame.paragraphs[0]
                run = para.runs[0] if para.runs else para.add_run()
                run.font.size = Pt(7)
                run.font.color.rgb = _DARK
                para.alignment = PP_ALIGN.CENTER

    def _fill_background(self, slide, color: RGBColor) -> None:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(
        self, slide, text: str, left, top, width, height,
        font_size=Pt(12), bold=False, color=_DARK, align=PP_ALIGN.LEFT,
    ) -> None:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color

    def _add_rect(self, slide, left, top, width, height, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()  # no border
