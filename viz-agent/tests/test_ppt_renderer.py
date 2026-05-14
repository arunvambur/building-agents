import os

import pytest

from core.dsl.schema import Aggregation, Axis, Chart, Filter, VisualizationSpec
from core.renderer.ppt.ppt import PPTRenderer

FILE_PREFIX = "file://"


def strip_prefix(result: str) -> str:
    assert result.startswith(FILE_PREFIX)
    return result[len(FILE_PREFIX):]


@pytest.fixture
def renderer():
    return PPTRenderer()


@pytest.fixture
def sample_data():
    return [
        {"town": "St Ives",  "rating": 4.8, "price_single": 140.0, "available_rooms": 6},
        {"town": "Newquay",  "rating": 4.5, "price_single": 120.0, "available_rooms": 5},
        {"town": "Falmouth", "rating": 4.2, "price_single": 95.0,  "available_rooms": 2},
        {"town": "Penzance", "rating": 4.7, "price_single": 130.0, "available_rooms": 3},
    ]


@pytest.fixture
def bar_spec():
    return VisualizationSpec(
        charts=[Chart(
            type="bar",
            x=Axis(field="town", type="dimension"),
            y=Axis(field="rating", type="measure"),
            aggregation=Aggregation(field="rating", op="avg"),
            title="Avg Rating by Town",
        )],
        output="ppt",
    )


def test_supports_ppt(renderer):
    assert renderer.supports("ppt") is True


def test_does_not_support_pdf(renderer):
    assert renderer.supports("pdf") is False


def test_render_returns_file_prefix(renderer, bar_spec, sample_data):
    result = renderer.render(bar_spec, sample_data)
    assert result.startswith(FILE_PREFIX)


def test_render_creates_pptx_file(renderer, bar_spec, sample_data):
    path = strip_prefix(renderer.render(bar_spec, sample_data))
    assert path.endswith(".pptx")
    assert os.path.exists(path)
    assert os.path.getsize(path) > 1000
    os.remove(path)


def test_render_correct_slide_count(renderer, bar_spec, sample_data):
    from pptx import Presentation
    path = strip_prefix(renderer.render(bar_spec, sample_data))
    prs = Presentation(path)
    # Title slide + 1 chart slide = 2
    assert len(prs.slides) == 2
    os.remove(path)


def test_render_with_empty_data(renderer, bar_spec):
    path = strip_prefix(renderer.render(bar_spec, []))
    assert os.path.exists(path)
    os.remove(path)


def test_render_multi_chart_slide_count(renderer, sample_data):
    spec = VisualizationSpec(
        charts=[
            Chart(type="bar", x=Axis(field="town", type="dimension"),
                  y=Axis(field="rating", type="measure"), title="Ratings"),
            Chart(type="line", x=Axis(field="town", type="dimension"),
                  y=Axis(field="price_single", type="measure"), title="Prices"),
        ],
        layout="grid",
        output="ppt",
    )
    from pptx import Presentation
    path = strip_prefix(renderer.render(spec, sample_data))
    prs = Presentation(path)
    # Title slide + 2 chart slides = 3
    assert len(prs.slides) == 3
    os.remove(path)


def test_render_with_filters_adds_filter_slide(renderer, sample_data):
    spec = VisualizationSpec(
        charts=[Chart(
            type="bar",
            x=Axis(field="town", type="dimension"),
            y=Axis(field="rating", type="measure"),
            title="Filtered",
        )],
        filters=[Filter(field="rating", op=">", value="4.3")],
        output="ppt",
    )
    from pptx import Presentation
    path = strip_prefix(renderer.render(spec, sample_data))
    prs = Presentation(path)
    # Title + chart + filters = 3
    assert len(prs.slides) == 3
    os.remove(path)


def test_render_pie_chart(renderer, sample_data):
    spec = VisualizationSpec(
        charts=[Chart(
            type="pie",
            x=Axis(field="town", type="dimension"),
            y=Axis(field="available_rooms", type="measure"),
            aggregation=Aggregation(field="available_rooms", op="sum"),
            title="Room Distribution",
        )],
        output="ppt",
    )
    path = strip_prefix(renderer.render(spec, sample_data))
    assert os.path.exists(path)
    os.remove(path)
